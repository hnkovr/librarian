"""librarian.docprops — deterministic document properties for catalog rows.

Per format:
  pptx — slides, pictures, code blocks (runs in monospace fonts), per-slide word
         stats, per-slide speaker-notes word stats, core properties (author, dates);
  pdf  — pages, per-page word stats, embedded image count;
  docx — paragraph word stats, inline images, core properties;
  md   — word count, fenced code blocks, image refs.
All heuristics are pure functions of file bytes + settings — same input, same row.
"""

from __future__ import annotations

import re
from pathlib import Path

from loguru import logger as log

from . import settings
from .utils import word_stats

_WORDS = None


def _words(text: str) -> int:
    global _WORDS
    if _WORDS is None:
        _WORDS = re.compile(settings.get("docprops", "words_split_regex"))
    return len([w for w in _WORDS.split(text.strip()) if w])


def _core_props(cp) -> dict:
    return {
        "author": cp.author or None,
        "last_modified_by": cp.last_modified_by or None,
        "created": cp.created.isoformat() if cp.created else None,
        "modified": cp.modified.isoformat() if cp.modified else None,
        "title": cp.title or None,
        "subject": cp.subject or None,
        "keywords": cp.keywords or None,
    }


def pptx_props(path: Path) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mono = {f.casefold() for f in settings.get("docprops", "monospace_fonts")}
    prs = Presentation(str(path))
    slide_words: list[int] = []
    notes_words: list[int] = []
    pictures = 0
    code_blocks = 0
    for slide in prs.slides:
        words = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if not shape.has_text_frame:
                continue
            has_mono = False
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    words += _words(run.text)
                    if run.font.name and run.font.name.casefold() in mono:
                        has_mono = True
            code_blocks += 1 if has_mono else 0
        slide_words.append(words)
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        notes_words.append(_words(notes))
    return {
        "format": "pptx",
        "slides": len(prs.slides),
        "pictures": pictures,
        "code_blocks": code_blocks,
        "words_per_slide": word_stats(slide_words),
        "notes_words_per_slide": word_stats(notes_words),
        **_core_props(prs.core_properties),
    }


def pdf_props(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    page_words: list[int] = []
    images = 0
    for page in reader.pages:
        try:
            page_words.append(_words(page.extract_text() or ""))
        except Exception as e:  # noqa: BLE001 — malformed page must not kill the sweep
            log.warning("pdf text extraction failed on a page of {}: {}", path.name, e)
            page_words.append(0)
        try:
            images += len(page.images)
        except Exception:  # noqa: BLE001
            pass
    meta = reader.metadata or {}
    return {
        "format": "pdf",
        "pages": len(reader.pages),
        "pictures": images,
        "words_per_page": word_stats(page_words),
        "author": (meta.get("/Author") or None) if meta else None,
        "created": str(meta.get("/CreationDate")) if meta and meta.get("/CreationDate") else None,
        "modified": str(meta.get("/ModDate")) if meta and meta.get("/ModDate") else None,
        "title": (meta.get("/Title") or None) if meta else None,
    }


def docx_props(path: Path) -> dict:
    import docx

    d = docx.Document(str(path))
    para_words = [_words(p.text) for p in d.paragraphs if p.text.strip()]
    return {
        "format": "docx",
        "paragraphs": len(para_words),
        "pictures": len(d.inline_shapes),
        "words_per_paragraph": word_stats(para_words),
        **_core_props(d.core_properties),
    }


def md_props(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace")
    fence = settings.get("docprops", "md_code_fence")
    return {
        "format": "md",
        "words": _words(re.sub(r"[#*_`>\[\]()!-]", " ", text)),
        "code_blocks": text.count(fence) // 2,
        "pictures": len(re.findall(r"!\[", text)),
    }


_HANDLERS = {"pptx": pptx_props, "pdf": pdf_props, "docx": docx_props, "md": md_props}


def props_for(path: Path) -> dict:
    """Dispatch by extension; unsupported formats get a minimal stat row."""
    handler = _HANDLERS.get(path.suffix.lower().lstrip("."))
    base = {"file": str(path), "bytes": path.stat().st_size}
    if handler is None:
        return {**base, "format": path.suffix.lower().lstrip(".") or "unknown"}
    try:
        return {**base, **handler(path)}
    except Exception as e:  # noqa: BLE001 — one broken file must not abort the catalog
        log.error("docprops failed for {}: {}", path, e)
        return {**base, "format": "error", "error": str(e)}
